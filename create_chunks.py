from langchain_community.document_loaders import YoutubeLoader, TextLoader, WebBaseLoader, PyPDFLoader
from langchain.text_splitter import RecursiveCharacterTextSplitter
import pandas as pd
from docx import Document
from pptx import Presentation
from dotenv import load_dotenv

load_dotenv()

# Function to create chunks from a YouTube video transcript
def chunk_youtube_video(url, chunk_size=1000, chunk_overlap=200):
    """
    Loads the transcript of a YouTube video and splits it into chunks.
    """
    loader = YoutubeLoader.from_youtube_url(url)
    transcript = loader.load()

    text_splitter = RecursiveCharacterTextSplitter(chunk_size=chunk_size, chunk_overlap=chunk_overlap)
    chunks = text_splitter.split_documents(transcript)

    return chunks

# Function to create chunks from a PDF file
def chunk_pdf(file_path, chunk_size=1000, chunk_overlap=200):
    """
    Reads a PDF file and creates chunks from its content.
    """
    loader = PyPDFLoader(file_path)
    docs = loader.load()

    text_splitter = RecursiveCharacterTextSplitter(chunk_size=chunk_size, chunk_overlap=chunk_overlap)
    chunks = text_splitter.split_documents(docs)

    return chunks

# Function to create chunks from an Excel file
def chunk_excel(file_path, chunk_size=1000, chunk_overlap=200):
    """
    Reads an Excel file and creates chunks from its content.
    """
    excel_data = pd.read_excel(file_path)
    text = "\n".join(excel_data.astype(str).apply(' '.join, axis=1))

    text_splitter = RecursiveCharacterTextSplitter(chunk_size=chunk_size, chunk_overlap=chunk_overlap)
    chunks = text_splitter.create_documents([text])

    return chunks

# Function to create chunks from a website URL
def chunk_website(url, chunk_size=1000, chunk_overlap=200):
    """
    Fetches the content of a webpage and splits it into chunks.
    """
    loader = WebBaseLoader(url)
    docs = loader.load()

    text_splitter = RecursiveCharacterTextSplitter(chunk_size=chunk_size, chunk_overlap=chunk_overlap)
    chunks = text_splitter.split_documents(docs)

    return chunks

# Function to create chunks from a plain text file
def chunk_txt_file(file_path, chunk_size=1000, chunk_overlap=200):
    """
    Reads a plain text file and creates chunks from its content.
    """
    loader = TextLoader(file_path)
    text = loader.load()

    text_splitter = RecursiveCharacterTextSplitter(chunk_size=chunk_size, chunk_overlap=chunk_overlap)
    chunks = text_splitter.split_documents(text)

    return chunks

# Function to create chunks from a Word document (.docx)
def chunk_word_doc(file_path, chunk_size=1000, chunk_overlap=200):
    """
    Reads a Word document (.docx) and creates chunks from its content.
    """
    doc = Document(file_path)
    text = "\n".join([para.text for para in doc.paragraphs])

    text_splitter = RecursiveCharacterTextSplitter(chunk_size=chunk_size, chunk_overlap=chunk_overlap)
    chunks = text_splitter.create_documents([text])

    return chunks

def chunk_pptx(file_path, chunk_size=1000, chunk_overlap=200):
    """
    Reads a PowerPoint (.pptx) file and creates chunks from its content.
    """
    # Load the PowerPoint presentation
    presentation = Presentation(file_path)
    text_content = []

    # Extract text from each slide
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_content.append(shape.text)

    full_text = "\n".join(text_content)

    text_splitter = RecursiveCharacterTextSplitter(chunk_size=chunk_size, chunk_overlap=chunk_overlap)
    chunks = text_splitter.create_documents([full_text])

    return chunks

# print(chunk_youtube_video("https://youtu.be/ABFqbY_rmEk?si=NL7eaE21XGBe9_Sp"))
# print(chunk_website("https://pytorch.org/get-started/locally/com"))
# print(chunk_pdf("data/data-pdf.pdf"))
# print(chunk_txt_file("data/data-txt.txt"))
# print(chunk_word_doc("data/data-docs.docx"))
# print(chunk_pptx("data/data-ppt.pptx"))
# print(chunk_excel("data/data-excel.xlsx"))