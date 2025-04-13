import pandas as pd
from docx import Document 
from pptx import Presentation
from dotenv import load_dotenv
from langchain_chroma import Chroma
from langchain.prompts import ChatPromptTemplate
from langchain.chains import create_retrieval_chain
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.chains.combine_documents import create_stuff_documents_chain
from langchain_google_genai import GoogleGenerativeAIEmbeddings, ChatGoogleGenerativeAI
from langchain_community.document_loaders import YoutubeLoader, TextLoader, WebBaseLoader, PyPDFLoader
from utils.constants import CHUNK_SIZE, CHUNK_OVERLAP, EMBEDDING_MODEL, LLM_MODEL

load_dotenv()

def chunk_youtube_video(url, chunk_size=CHUNK_SIZE, chunk_overlap=CHUNK_OVERLAP):
    """
    Loads the transcript of a YouTube video and splits it into chunks.
    """
    loader = YoutubeLoader.from_youtube_url(url)
    transcript = loader.load()

    text_splitter = RecursiveCharacterTextSplitter(chunk_size=chunk_size, chunk_overlap=chunk_overlap)
    chunks = text_splitter.split_documents(transcript)

    return chunks

def chunk_pdf(file_path, chunk_size=CHUNK_SIZE, chunk_overlap=CHUNK_OVERLAP):
    """
    Reads a PDF file and creates chunks from its content.
    """
    loader = PyPDFLoader(file_path)
    docs = loader.load()

    text_splitter = RecursiveCharacterTextSplitter(chunk_size=chunk_size, chunk_overlap=chunk_overlap)
    chunks = text_splitter.split_documents(docs)

    return chunks

def chunk_excel(file_path, chunk_size=CHUNK_SIZE, chunk_overlap=CHUNK_OVERLAP):
    """
    Reads an Excel file and creates chunks from its content.
    """
    excel_data = pd.read_excel(file_path)
    text = "\n".join(excel_data.astype(str).apply(' '.join, axis=1))

    text_splitter = RecursiveCharacterTextSplitter(chunk_size=chunk_size, chunk_overlap=chunk_overlap)
    chunks = text_splitter.create_documents([text])

    return chunks

def chunk_website(url, chunk_size=CHUNK_SIZE, chunk_overlap=CHUNK_OVERLAP):
    """
    Fetches the content of a webpage and splits it into chunks.
    """
    loader = WebBaseLoader(url)
    docs = loader.load()

    text_splitter = RecursiveCharacterTextSplitter(chunk_size=chunk_size, chunk_overlap=chunk_overlap)
    chunks = text_splitter.split_documents(docs)

    return chunks

def chunk_txt_file(file_path, chunk_size=CHUNK_SIZE, chunk_overlap=CHUNK_OVERLAP):
    """
    Reads a plain text file and creates chunks from its content.
    """
    loader = TextLoader(file_path)
    text = loader.load()

    text_splitter = RecursiveCharacterTextSplitter(chunk_size=chunk_size, chunk_overlap=chunk_overlap)
    chunks = text_splitter.split_documents(text)

    return chunks

def chunk_word_doc(file_path, chunk_size=CHUNK_SIZE, chunk_overlap=CHUNK_OVERLAP):
    """
    Reads a Word document (.docx) and creates chunks from its content.
    """
    doc = Document(file_path)
    text = "\n".join([para.text for para in doc.paragraphs])

    text_splitter = RecursiveCharacterTextSplitter(chunk_size=chunk_size, chunk_overlap=chunk_overlap)
    chunks = text_splitter.create_documents([text])

    return chunks

def chunk_pptx(file_path, chunk_size=CHUNK_SIZE, chunk_overlap=CHUNK_OVERLAP):
    """
    Reads a PowerPoint (.pptx) file and creates chunks from its content.
    """
    # Load the PowerPoint presentation
    presentation = Presentation(file_path)
    text_content = []

    # Extract text from each slide
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_content.append(shape.text)

    full_text = "\n".join(text_content)

    text_splitter = RecursiveCharacterTextSplitter(chunk_size=chunk_size, chunk_overlap=chunk_overlap)
    chunks = text_splitter.create_documents([full_text])

    return chunks

def create_vectorstore(chunks, persist_directory):
    """
    Creates a vectorstore from the chunks.
    """
    return Chroma.from_documents(documents=chunks, embedding=GoogleGenerativeAIEmbeddings(model=EMBEDDING_MODEL), persist_directory=persist_directory)
    # return Chroma.from_documents(documents=chunks, embedding=GoogleGenerativeAIEmbeddings(model=EMBEDDING_MODEL))


def create_rag_chain(vectorstore):
    """
    Creates a retrieval-augmented generation (RAG) chain.
    """

    retriever = vectorstore.as_retriever(search_type="similarity", search_kwargs={"k": 10}) # k: 10 means retrive ten similar documents from 96 documents in the vectorstore

    llm = ChatGoogleGenerativeAI(model=LLM_MODEL, temperature=0.3, max_tokens=None)

    system_prompt = (
            "You are an assistant for question-answering tasks. "
            "Use the following pieces of retrieved context to answer "
            "the question. If you don't know the answer, say that you don't know."
        )

    prompt = ChatPromptTemplate.from_messages([
        ("system", system_prompt),
        ("user", "Answer the question based on the context below:\n\n{context}\n\nQuestion: {input}"),
    ])

    question_answer_chain =  create_stuff_documents_chain(llm, prompt)
    rag_chain = create_retrieval_chain(retriever, question_answer_chain)

    return rag_chain

def chat_with_rag_chain(rag_chain, question):
    """
    Interacts with the RAG chain to get an answer for a given question.
    """
    response = rag_chain.invoke({"input": question})
    return response["answer"]

def main():
    """
    Main function to run the RAG chain.
    """
    
    chunks = chunk_website("https://github.com/sanket-164")
    # chunks = chunk_youtube_video("https://youtu.be/QhMO5SSmiaA?si=mhgXeY5AdQRvuYAl")
    # chunks = chunk_pdf("data/data-pdf.pdf")
    # chunks = chunk_txt_file("data/data-txt.txt")
    # chunks = chunk_word_doc("data/data-docs.docx")
    # chunks = chunk_pptx("data/data-ppt.pptx")
    # chunks = chunk_excel("data/data-excel.xlsx")

    vectorstore = create_vectorstore(chunks, "vectorstore")

    rag_chain = create_rag_chain(vectorstore)

    while True:
        question = input("You: ")
        
        if question.lower() == "exit" or question.lower() == "quit" or question.lower() == "q":
            break

        response = chat_with_rag_chain(rag_chain, question)
        
        print("AI:", response)

    return

if __name__ == "__main__":
    main()